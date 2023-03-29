import typing


import pptx
from pptx.presentation import Presentation
from pptx.slide import Slide, SlideLayout

from slidesgen.presentation_generation.data import Template
from slidesgen.content_generation.data import Content


class PresentationGenerator:

	def __init__(self, template: typing.Optional[Template] = None):
		self.__template = template
		if template is None:
			self.__template = Template(None)

	def __get_layout(self, presentation: Presentation, type_: int) -> SlideLayout:
		return presentation.slide_layouts[self.__template.get_layout(type_)]

	def __create_title_slide(self, presentation: Presentation, title: str):
		layout = self.__get_layout(presentation, Template.Layouts.TITLE)

		slide: Slide = presentation.slides.add_slide(layout)
		slide.placeholders[0].text = title

	def __create_content_slide(self, presentation: Presentation, content: Content):
		layout = self.__get_layout(presentation, Template.Layouts.TITLE_AND_CONTENT)

		slide = presentation.slides.add_slide(layout)
		title_idx, body_idx = [shape.placeholder_format.idx for shape in slide.placeholders]

		if content.title is not None:
			title_shape = slide.placeholders[title_idx]
			title_shape.text = content.title

		body_shape = slide.placeholders[body_idx]

		for point in content.body:
			paragraph = body_shape.text_frame.add_paragraph()
			paragraph.text = point

	def __create_end_slide(self, presentation: Presentation):
		layout = self.__get_layout(presentation, Template.Layouts.END)
		slide = presentation.slides.add_slide(layout)

	def __create_presentation(self) -> Presentation:
		presentation = pptx.Presentation(self.__template.path)
		return presentation

	def generate(
			self,
			title: str,
			contents: typing.List[Content],
			out_path: str,
	):

		presentation = self.__create_presentation()
		self.__create_title_slide(presentation, title)

		for content in contents:
			self.__create_content_slide(presentation, content)

		self.__create_end_slide(presentation)

		presentation.save(out_path)
