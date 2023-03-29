import typing

import pptx
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from slidesgen.content_generation.data import Content
from .file_content_extractor import FileContentExtractor


class PPTXExtractor(FileContentExtractor):

	@staticmethod
	def __extract_body(slide: Slide) -> typing.List[str]:
		return max(([
			shape.text
			for shape in slide.shapes
			if isinstance(shape, SlidePlaceholder)
		]), key=lambda text: len(text)).split("\n")

	@staticmethod
	def __extract_title(slide: Slide) -> typing.Optional[str]:
		for placeholder in slide.placeholders:
			if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE:
				return placeholder.text
		return None

	@staticmethod
	def __extract_slide(slide: Slide) -> typing.List[str]:
		return max(([
			shape.text.split("\n")
			for shape in slide.shapes
			if isinstance(shape, SlidePlaceholder)
		]), key=lambda text: len(text))

	def extract(self, file_path: str, *args, **kwargs) -> typing.List[Content]:
		presentation: Presentation = pptx.Presentation(file_path)
		return [
			Content(
				title=self.__extract_title(slide),
				body=self.__extract_body(slide)
			)
			for slide in presentation.slides
		]