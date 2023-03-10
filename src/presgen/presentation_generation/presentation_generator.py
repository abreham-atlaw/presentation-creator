import typing

from pptx import Presentation


class PresentationGenerator:

	def __init__(self):
		pass

	def __generate_slide(self, presentation: Presentation, contents: typing.List[str]):
		layout = presentation.slide_layouts[1]
		slide = presentation.slides.add_slide(layout)
		body_shape = slide.shapes.placeholders[1]

		for point in contents:
			paragraph = body_shape.text_frame.add_paragraph()
			paragraph.text = point

	def generate(self, contents: typing.List[typing.List[str]], out_path: str):
		presentation = Presentation()
		for content in contents:
			self.__generate_slide(presentation, content)
		presentation.save(out_path)
